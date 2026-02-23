"""
file_to_structured_all_md_hardened.py

File-agnostic context ingestion script (hardened).

Outputs (per file):
  - <stem>_structured.json
  - <stem>_structured.md

Supported:
  - .pptx
  - .docx
  - .xlsx / .xlsm
  - .txt / .md
  - .pdf

Key hardening improvements vs prior version:
  - DOCX: preserves true document order (paragraphs + tables interleaved)
  - PPTX: extracts speaker notes (if present)
  - Optional: per-file envelopes and per-section hashes for auditability

Dependencies:
  pip install python-pptx python-docx openpyxl pdfplumber
"""

import os
import re
import json
import hashlib
from typing import Any, Dict, List, Optional, Iterable, Tuple

# -------------------------------------------------------------------
# CONFIG
# -------------------------------------------------------------------
OUTPUT_DIR_NAME = "output"
INPUT_DIR_NAME = ""  # "" = same folder as script

INCLUDE_EXTENSIONS = {
    ".pptx",
    ".docx",
    ".xlsx",
    ".xlsm",
    ".txt",
    ".md",
    ".pdf",
}

# Density thresholds (tweak freely)
LOW_DENSITY_WORDS = 80
HIGH_DENSITY_WORDS = 300

# Optional "lossless" affordances
INCLUDE_FILE_ENVELOPES_IN_MD = True  # wraps each structured MD in BEGIN/END markers
INCLUDE_SECTION_HASHES = True        # adds sha256 checksum per section (JSON + MD)
INCLUDE_SECTION_LINE_COUNTS = True   # adds line_count per section (JSON + MD)

# Processing behavior
WALK_SUBDIRS = False  # True = recursively walk input_dir; False = only top-level files
SKIP_PREFIXES = ("~$",)  # Office temp files etc.

# -------------------------------------------------------------------
# Shared helpers
# -------------------------------------------------------------------
def slugify(text: str, default: str = "section") -> str:
    if not text:
        return default
    slug = text.lower()
    slug = re.sub(r"[^a-z0-9]+", "_", slug)
    slug = re.sub(r"_+", "_", slug)
    return slug.strip("_") or default


def write_json(data: Dict[str, Any], path: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def normalize_cell_for_md(text: Any) -> str:
    if text is None:
        return ""
    # Keep line breaks, but normalize into a single cell using " / " separators.
    parts = [p.strip() for p in str(text).splitlines() if p.strip()]
    return " / ".join(parts)


def render_table_markdown_from_matrix(table_rows: List[List[Any]]) -> str:
    if not table_rows:
        return ""

    max_cols = max(len(r) for r in table_rows)
    rows = [r + [""] * (max_cols - len(r)) for r in table_rows]

    header = rows[0]
    body = rows[1:]

    lines = []
    lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")

    for row in body:
        lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in row) + " |")

    return "\n".join(lines)


def compute_density(body_blocks: List[str], table_count: int) -> Dict[str, Any]:
    words = sum(len(b.split()) for b in body_blocks)
    if words < LOW_DENSITY_WORDS and table_count == 0:
        density = "low"
    elif words > HIGH_DENSITY_WORDS or table_count >= 2:
        density = "high"
    else:
        density = "medium"

    return {
        "word_count": words,
        "table_count": table_count,
        "density": density,
    }


def build_table_object(
    section_slug: str,
    table_index: int,
    matrix: List[List[str]],
) -> Optional[Dict[str, Any]]:
    if not matrix or not any(any((c or "").strip() for c in r) for r in matrix):
        return None

    max_cols = max(len(r) for r in matrix)
    norm = [r + [""] * (max_cols - len(r)) for r in matrix]

    # If first row has any content, treat as header; else treat as headerless table
    if any((c or "").strip() for c in norm[0]):
        headers = norm[0]
        rows = norm[1:]
    else:
        headers = []
        rows = norm

    return {
        "id": f"{section_slug}_table_{table_index}",
        "headers": headers,
        "rows": rows,
        "raw": norm,
    }


def sha256_text(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8", errors="ignore")).hexdigest()


def compute_section_checksum(sec: Dict[str, Any]) -> str:
    # Stable-ish checksum over content-bearing fields
    payload = {
        "title": sec.get("title") or "",
        "layout": sec.get("layout") or "",
        "body_blocks": sec.get("body_blocks") or [],
        "tables": sec.get("tables") or [],
        "notes": sec.get("notes") or "",
        "source": sec.get("source") or {},
    }
    return sha256_text(json.dumps(payload, ensure_ascii=False, sort_keys=True))


def compute_line_count(body_blocks: List[str], notes: str = "") -> int:
    lines = 0
    for b in body_blocks or []:
        lines += len([ln for ln in b.splitlines() if ln.strip()])
    if notes:
        lines += len([ln for ln in notes.splitlines() if ln.strip()])
    return lines


# -------------------------------------------------------------------
# Markdown writer
# -------------------------------------------------------------------
def write_markdown(doc: Dict[str, Any], path: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        if INCLUDE_FILE_ENVELOPES_IN_MD:
            f.write(f"<!-- BEGIN_SOURCE_FILE: {doc['source_file']} -->\n\n")

        f.write(f"# Source: {doc['source_file']}\n")
        f.write(f"_File type: {doc['file_type']}_\n\n")
        f.write("---\n\n")

        for sec in doc["sections"]:
            idx = sec["index"]
            title = sec["title"] or f"Section {idx}"

            f.write(f"## SECTION {idx}: {title}\n\n")
            f.write(f"_Slug: {sec['slug']}_\n\n")
            f.write(f"_Layout: {sec['layout']}_\n\n")

            # Trace metadata
            src = sec["source"]
            f.write(f"_Source: {src['file']} | {src['type']} | {src['location']}_\n\n")

            # Density
            sig = sec["signals"]
            f.write(
                f"_Signals: {sig['density']} "
                f"(words={sig['word_count']}, tables={sig['table_count']})_\n\n"
            )

            if INCLUDE_SECTION_LINE_COUNTS and "line_count" in sig:
                f.write(f"_Line count: {sig['line_count']}_\n\n")

            if INCLUDE_SECTION_HASHES and "sha256" in sig:
                f.write(f"_Section sha256: {sig['sha256']}_\n\n")

            # Body
            for block in sec["body_blocks"]:
                for line in block.splitlines():
                    if line.strip():
                        f.write(f"- {line.strip()}\n")
                f.write("\n")

            # Tables
            for i, table in enumerate(sec["tables"], start=1):
                f.write(f"**Table {i} ({table['id']}):**\n\n")
                matrix = (
                    [table["headers"]] + table["rows"]
                    if table.get("headers")
                    else table.get("raw", [])
                )
                f.write(render_table_markdown_from_matrix(matrix) + "\n\n")

            # Notes
            if sec.get("notes"):
                f.write("**Notes:**\n\n")
                f.write(sec["notes"].strip() + "\n\n")

            f.write("---\n\n")

        if INCLUDE_FILE_ENVELOPES_IN_MD:
            f.write(f"<!-- END_SOURCE_FILE: {doc['source_file']} -->\n")


# -------------------------------------------------------------------
# PPTX extractor
# -------------------------------------------------------------------
def extract_pptx(path: str) -> Dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(s.shapes)
            else:
                yield s

    def text(shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()

    def table(shape):
        if not getattr(shape, "has_table", False):
            return None
        return [
            [
                "\n".join(p.text for p in cell.text_frame.paragraphs).strip()
                if getattr(cell, "text_frame", None)
                else ""
                for cell in row.cells
            ]
            for row in shape.table.rows
        ]

    def notes(slide) -> str:
        # Speaker notes extraction (hardened)
        try:
            if not slide.has_notes_slide:
                return ""
            ns = slide.notes_slide
            if not ns or not ns.notes_text_frame:
                return ""
            txt = "\n".join(p.text for p in ns.notes_text_frame.paragraphs).strip()
            return txt
        except Exception:
            return ""

    prs = Presentation(path)
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        title = text(slide.shapes.title) if slide.shapes.title else ""
        slug = slugify(title, f"section_{i}")

        body_blocks: List[str] = []
        raw_tables: List[List[List[str]]] = []

        for shape in walk(slide.shapes):
            if slide.shapes.title and shape == slide.shapes.title:
                continue

            t = table(shape)
            if t:
                raw_tables.append(t)
                continue

            txt = text(shape)
            if txt:
                body_blocks.append(txt)

        tables = []
        for ti, matrix in enumerate(raw_tables, start=1):
            obj = build_table_object(slug, ti, matrix)
            if obj:
                tables.append(obj)

        sec_notes = notes(slide)
        signals = compute_density(body_blocks, len(tables))
        if INCLUDE_SECTION_LINE_COUNTS:
            signals["line_count"] = compute_line_count(body_blocks, sec_notes)
        sec = {
            "index": i,
            "slug": slug,
            "label": f"Slide {i}",
            "layout": slide.slide_layout.name if slide.slide_layout else "",
            "title": title,
            "body_blocks": body_blocks,
            "tables": tables,
            "notes": sec_notes,
            "source": {
                "file": os.path.basename(path),
                "type": "pptx",
                "location": f"Slide {i}",
            },
            "signals": signals,
        }
        if INCLUDE_SECTION_HASHES:
            sec["signals"]["sha256"] = compute_section_checksum(sec)

        sections.append(sec)

    return {
        "source_file": os.path.basename(path),
        "file_type": "pptx",
        "sections": sections,
    }


# -------------------------------------------------------------------
# DOCX extractor (ordered paragraphs + tables)
# -------------------------------------------------------------------
def extract_docx(path: str) -> Dict[str, Any]:
    from docx import Document
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = Document(path)
    sections: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    idx = 0
    table_counters: Dict[str, int] = {}  # per-section table counter

    def start(title: str, style: str) -> Dict[str, Any]:
        nonlocal idx, current
        idx += 1
        sec = {
            "index": idx,
            "slug": slugify(title, f"section_{idx}"),
            "label": f"Section {idx}",
            "layout": style or "Document",
            "title": title,
            "body_blocks": [],
            "tables": [],
            "notes": "",
            "source": {
                "file": os.path.basename(path),
                "type": "docx",
                "location": f"Section {idx}",
            },
            # signals added later
        }
        sections.append(sec)
        current = sec
        table_counters[sec["slug"]] = 0
        return sec

    def get_or_start_default() -> Dict[str, Any]:
        nonlocal current
        if current:
            return current
        stem = os.path.splitext(os.path.basename(path))[0]
        return start(stem, "Document")

    def is_heading(par: Paragraph) -> Tuple[bool, str]:
        try:
            style = par.style.name if par.style else ""
        except Exception:
            style = ""
        txt = (par.text or "").strip()
        if txt and style.startswith("Heading"):
            return True, style
        return False, style

    def table_to_matrix(tbl: Table) -> List[List[str]]:
        matrix = []
        for row in tbl.rows:
            matrix.append([ (cell.text or "").strip() for cell in row.cells ])
        return matrix

    # Iterate body children to preserve order
    body_elm = doc.element.body  # type: ignore[attr-defined]
    for child in body_elm.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]  # strip namespace
        if tag == "p":
            par = Paragraph(child, doc)
            txt = (par.text or "").strip()
            is_h, style = is_heading(par)
            if is_h:
                start(txt, style)
            else:
                sec = get_or_start_default()
                if txt:
                    sec["body_blocks"].append(txt)
        elif tag == "tbl":
            tbl = Table(child, doc)
            sec = get_or_start_default()
            slug = sec["slug"]
            table_counters[slug] = table_counters.get(slug, 0) + 1
            matrix = table_to_matrix(tbl)
            obj = build_table_object(slug, table_counters[slug], matrix)
            if obj:
                sec["tables"].append(obj)
        else:
            # Other element types are ignored (e.g., sectPr)
            continue

    # Signals + optional hashes
    for sec in sections:
        sec_notes = sec.get("notes") or ""
        sig = compute_density(sec["body_blocks"], len(sec["tables"]))
        if INCLUDE_SECTION_LINE_COUNTS:
            sig["line_count"] = compute_line_count(sec["body_blocks"], sec_notes)
        sec["signals"] = sig
        if INCLUDE_SECTION_HASHES:
            sec["signals"]["sha256"] = compute_section_checksum(sec)

    return {
        "source_file": os.path.basename(path),
        "file_type": "docx",
        "sections": sections,
    }


# -------------------------------------------------------------------
# XLSX extractor
# -------------------------------------------------------------------
def extract_xlsx(path: str) -> Dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    sections = []

    for i, sheet in enumerate(wb.worksheets, start=1):
        matrix: List[List[str]] = []
        for r in sheet.iter_rows():
            matrix.append(["" if c.value is None else str(c.value).strip() for c in r])

        slug = slugify(sheet.title, f"sheet_{i}")
        table = build_table_object(slug, 1, matrix)

        body_blocks: List[str] = []
        tables = [table] if table else []

        sig = compute_density(body_blocks, len(tables))
        if INCLUDE_SECTION_LINE_COUNTS:
            sig["line_count"] = compute_line_count(body_blocks, "")
        sec = {
            "index": i,
            "slug": slug,
            "label": f"Sheet: {sheet.title}",
            "layout": "Spreadsheet",
            "title": sheet.title,
            "body_blocks": body_blocks,
            "tables": tables,
            "notes": "",
            "source": {
                "file": os.path.basename(path),
                "type": "xlsx",
                "location": f"Sheet: {sheet.title}",
            },
            "signals": sig,
        }
        if INCLUDE_SECTION_HASHES:
            sec["signals"]["sha256"] = compute_section_checksum(sec)
        sections.append(sec)

    return {
        "source_file": os.path.basename(path),
        "file_type": "xlsx",
        "sections": sections,
    }


# -------------------------------------------------------------------
# PDF extractor
# -------------------------------------------------------------------
def extract_pdf(path: str) -> Dict[str, Any]:
    import pdfplumber

    sections = []

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            tables_raw = page.extract_tables() or []

            body_blocks: List[str] = []
            if text.strip():
                cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
                body_blocks.append(cleaned)

            tables = []
            slug = slugify(f"page_{i}", f"page_{i}")

            for ti, matrix in enumerate(tables_raw, start=1):
                norm_matrix = [[("" if c is None else str(c).strip()) for c in row] for row in matrix]
                obj = build_table_object(slug, ti, norm_matrix)
                if obj:
                    tables.append(obj)

            sig = compute_density(body_blocks, len(tables))
            if INCLUDE_SECTION_LINE_COUNTS:
                sig["line_count"] = compute_line_count(body_blocks, "")
            sec = {
                "index": i,
                "slug": slug,
                "label": f"Page {i}",
                "layout": "PDF Page",
                "title": f"Page {i}",
                "body_blocks": body_blocks,
                "tables": tables,
                "notes": "",
                "source": {
                    "file": os.path.basename(path),
                    "type": "pdf",
                    "location": f"Page {i}",
                },
                "signals": sig,
            }
            if INCLUDE_SECTION_HASHES:
                sec["signals"]["sha256"] = compute_section_checksum(sec)
            sections.append(sec)

    return {
        "source_file": os.path.basename(path),
        "file_type": "pdf",
        "sections": sections,
    }


# -------------------------------------------------------------------
# TXT / MD extractor
# -------------------------------------------------------------------
def extract_text_file(path: str, ftype: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read().strip()

    body = [content] if content else []
    sig = compute_density(body, 0)
    if INCLUDE_SECTION_LINE_COUNTS:
        sig["line_count"] = compute_line_count(body, "")
    sec = {
        "index": 1,
        "slug": slugify(os.path.splitext(os.path.basename(path))[0]),
        "label": "Section 1",
        "layout": "Text",
        "title": os.path.splitext(os.path.basename(path))[0],
        "body_blocks": body,
        "tables": [],
        "notes": "",
        "source": {
            "file": os.path.basename(path),
            "type": ftype,
            "location": "Entire document",
        },
        "signals": sig,
    }
    if INCLUDE_SECTION_HASHES:
        sec["signals"]["sha256"] = compute_section_checksum(sec)

    return {
        "source_file": os.path.basename(path),
        "file_type": ftype,
        "sections": [sec],
    }


# -------------------------------------------------------------------
# Dispatcher
# -------------------------------------------------------------------
def extract_file(path: str) -> Dict[str, Any]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext in (".xlsx", ".xlsm"):
        return extract_xlsx(path)
    if ext == ".txt":
        return extract_text_file(path, "txt")
    if ext == ".md":
        return extract_text_file(path, "md")
    if ext == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported file: {ext}")


# -------------------------------------------------------------------
# File discovery
# -------------------------------------------------------------------
def iter_input_files(input_dir: str) -> Iterable[str]:
    if WALK_SUBDIRS:
        for root, _, files in os.walk(input_dir):
            for fn in files:
                if fn.startswith(SKIP_PREFIXES):
                    continue
                if os.path.splitext(fn)[1].lower() in INCLUDE_EXTENSIONS:
                    yield os.path.join(root, fn)
    else:
        for fn in os.listdir(input_dir):
            if fn.startswith(SKIP_PREFIXES):
                continue
            if os.path.splitext(fn)[1].lower() in INCLUDE_EXTENSIONS:
                yield os.path.join(input_dir, fn)


# -------------------------------------------------------------------
# Main
# -------------------------------------------------------------------
def main():
    base = os.path.dirname(os.path.abspath(__file__))
    input_dir = os.path.join(base, INPUT_DIR_NAME) if INPUT_DIR_NAME else base
    output_dir = os.path.join(base, OUTPUT_DIR_NAME)
    os.makedirs(output_dir, exist_ok=True)

    files = list(iter_input_files(input_dir))
    if not files:
        print("No supported files found.")
        return

    failures: List[Tuple[str, str]] = []

    for path in files:
        fn = os.path.basename(path)
        print(f"Processing: {fn}")
        try:
            doc = extract_file(path)

            stem = os.path.splitext(fn)[0]
            write_json(doc, os.path.join(output_dir, f"{stem}_structured.json"))
            write_markdown(doc, os.path.join(output_dir, f"{stem}_structured.md"))
        except Exception as e:
            failures.append((fn, str(e)))

    if failures:
        print("\nCompleted with errors:")
        for fn, err in failures:
            print(f" - {fn}: {err}")
    else:
        print("\nDone.")


if __name__ == "__main__":
    main()
