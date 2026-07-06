"""
file_to_structured_all_md.py

File-agnostic context ingestion script.
Recursively processes all supported files under --input.

Outputs (per file):
  - <stem>_structured.json
  - <stem>_structured.md

Supported:
  - .pptx   (text, tables, speaker notes)
  - .docx   (text, tables interleaved in document order)
  - .xlsx / .xlsm
  - .txt / .md
  - .pdf    (text + vector tables; OCR fallback for image-heavy pages)

Dependencies (core):
  pip install python-pptx python-docx openpyxl pdfplumber

Dependencies (OCR fallback — optional):
  pip install pdf2image pytesseract
  + Tesseract binary: https://github.com/UB-Mannheim/tesseract/wiki

USAGE
  # Explicit paths
  python /path/to/scripts/file_to_structured_all_md.py \\
      --input /path/to/project \\
      --output /path/to/project/output

  # Defaults: input = cwd, output = <input>/output, timestamped subfolder
  python /path/to/scripts/file_to_structured_all_md.py

  # Skip timestamp subfolder
  python /path/to/scripts/file_to_structured_all_md.py --no-timestamp
"""

import os
import re
import json
import argparse
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


# ---------------------------------------------------------------------------
# CONFIG
# ---------------------------------------------------------------------------

INCLUDE_EXTENSIONS = {
    ".pptx",
    ".docx",
    ".xlsx",
    ".xlsm",
    ".txt",
    ".md",
    ".pdf",
}

# Skip Office temp files and already-converted outputs
SKIP_PREFIXES = ("~$",)
SKIP_SUFFIXES = ("_structured.md", "_structured.json")

# Density thresholds
LOW_DENSITY_WORDS  = 80
HIGH_DENSITY_WORDS = 300


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def slugify(text: str, default: str = "section") -> str:
    if not text:
        return default
    slug = text.lower()
    slug = re.sub(r"[^a-z0-9]+", "_", slug)
    slug = re.sub(r"_+", "_", slug)
    return slug.strip("_") or default


def write_json(data: Dict[str, Any], path: Path) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def normalize_cell_for_md(text: Any) -> str:
    if text is None:
        return ""
    parts = [p.strip() for p in str(text).splitlines() if p.strip()]
    return " / ".join(parts)


def render_table_markdown_from_matrix(table_rows: List[List[Any]]) -> str:
    if not table_rows:
        return ""

    max_cols = max(len(r) for r in table_rows)
    rows = [r + [""] * (max_cols - len(r)) for r in table_rows]

    header = rows[0]
    body   = rows[1:]

    lines = []
    lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in body:
        lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in row) + " |")

    return "\n".join(lines)


def compute_density(body_blocks: List[str], table_count: int) -> Dict[str, Any]:
    words = sum(len(b.split()) for b in body_blocks)
    if words < LOW_DENSITY_WORDS and table_count == 0:
        density = "low"
    elif words > HIGH_DENSITY_WORDS or table_count >= 2:
        density = "high"
    else:
        density = "medium"
    return {"word_count": words, "table_count": table_count, "density": density}


def build_table_object(
    section_slug: str,
    table_index: int,
    matrix: List[List[str]],
) -> Optional[Dict[str, Any]]:
    if not matrix or not any(any(c.strip() for c in r) for r in matrix):
        return None

    max_cols = max(len(r) for r in matrix)
    norm = [r + [""] * (max_cols - len(r)) for r in matrix]

    if any(c.strip() for c in norm[0]):
        headers = norm[0]
        rows    = norm[1:]
    else:
        headers = []
        rows    = norm

    return {
        "id":      f"{section_slug}_table_{table_index}",
        "headers": headers,
        "rows":    rows,
        "raw":     norm,
    }


# ---------------------------------------------------------------------------
# Markdown writer
# ---------------------------------------------------------------------------

def write_markdown(doc: Dict[str, Any], path: Path) -> None:
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"# Source: {doc['source_file']}\n")
        f.write(f"_File type: {doc['file_type']}_\n\n")
        f.write("---\n\n")

        for sec in doc["sections"]:
            idx   = sec["index"]
            title = sec["title"] or f"Section {idx}"

            f.write(f"## SECTION {idx}: {title}\n\n")
            f.write(f"_Slug: {sec['slug']}_\n\n")
            f.write(f"_Layout: {sec['layout']}_\n\n")

            src = sec["source"]
            f.write(f"_Source: {src['file']} | {src['type']} | {src['location']}_\n\n")

            sig = sec["signals"]
            f.write(
                f"_Signals: {sig['density']} "
                f"(words={sig['word_count']}, tables={sig['table_count']})_\n\n"
            )

            for block in sec["body_blocks"]:
                for line in block.splitlines():
                    if line.strip():
                        f.write(f"- {line.strip()}\n")
                f.write("\n")

            for i, table in enumerate(sec["tables"], start=1):
                f.write(f"**Table {i} ({table['id']}):**\n\n")
                matrix = (
                    [table["headers"]] + table["rows"]
                    if table["headers"]
                    else table["raw"]
                )
                f.write(render_table_markdown_from_matrix(matrix) + "\n\n")

            if sec.get("notes"):
                f.write("**Notes:**\n\n")
                f.write(sec["notes"] + "\n\n")

            f.write("---\n\n")


# ---------------------------------------------------------------------------
# PPTX extractor
# FIX: speaker notes are now captured (were always empty string before)
# ---------------------------------------------------------------------------

def extract_pptx(path: Path) -> Dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(s.shapes)
            else:
                yield s

    def get_text(shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()

    def get_table(shape) -> Optional[List[List[str]]]:
        if not getattr(shape, "has_table", False):
            return None
        return [
            [
                "\n".join(p.text for p in cell.text_frame.paragraphs).strip()
                if cell.text_frame else ""
                for cell in row.cells
            ]
            for row in shape.table.rows
        ]

    def get_notes(slide) -> str:
        """Extract speaker notes text from a slide."""
        try:
            notes_slide = slide.notes_slide
            if not notes_slide:
                return ""
            tf = notes_slide.notes_text_frame
            if not tf:
                return ""
            return "\n".join(p.text for p in tf.paragraphs).strip()
        except Exception:
            return ""

    prs = Presentation(str(path))
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        title = get_text(slide.shapes.title) if slide.shapes.title else ""
        slug  = slugify(title, f"section_{i}")

        body_blocks = []
        raw_tables  = []

        for shape in walk(slide.shapes):
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            t = get_table(shape)
            if t is not None:
                raw_tables.append(t)
                continue
            txt = get_text(shape)
            if txt:
                body_blocks.append(txt)

        tables = []
        for ti, matrix in enumerate(raw_tables, start=1):
            obj = build_table_object(slug, ti, matrix)
            if obj:
                tables.append(obj)

        notes   = get_notes(slide)
        signals = compute_density(body_blocks, len(tables))

        sections.append({
            "index":        i,
            "slug":         slug,
            "label":        f"Slide {i}",
            "layout":       slide.slide_layout.name if slide.slide_layout else "",
            "title":        title,
            "body_blocks":  body_blocks,
            "tables":       tables,
            "notes":        notes,
            "source": {
                "file":     path.name,
                "type":     "pptx",
                "location": f"Slide {i}",
            },
            "signals": signals,
        })

    return {
        "source_file": path.name,
        "file_type":   "pptx",
        "sections":    sections,
    }


# ---------------------------------------------------------------------------
# DOCX extractor
# FIX: tables are now placed in the section where they appear in document
#      order, not all dumped onto the last section.
# ---------------------------------------------------------------------------

def extract_docx(path: Path) -> Dict[str, Any]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph as DocxParagraph
    from docx.table import Table as DocxTable

    doc = Document(str(path))
    sections: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    idx = 0
    table_counters: Dict[str, int] = {}  # slug → running table index

    def start_section(title: str, style: str) -> Dict[str, Any]:
        nonlocal idx
        idx += 1
        sec = {
            "index":       idx,
            "slug":        slugify(title, f"section_{idx}"),
            "label":       f"Section {idx}",
            "layout":      style,
            "title":       title,
            "body_blocks": [],
            "tables":      [],
            "notes":       "",
            "source": {
                "file":     path.name,
                "type":     "docx",
                "location": f"Section {idx}",
            },
        }
        sections.append(sec)
        return sec

    # Iterate body children in document order so tables land in the right section
    for child in doc.element.body:
        tag = child.tag

        if tag == qn("w:p"):
            para  = DocxParagraph(child, doc)
            txt   = para.text.strip()
            style = para.style.name if para.style else ""

            if style.startswith("Heading") and txt:
                current = start_section(txt, style)
            else:
                if not current:
                    current = start_section(
                        path.stem, "Document"
                    )
                if txt:
                    current["body_blocks"].append(txt)

        elif tag == qn("w:tbl"):
            if not current:
                current = start_section(path.stem, "Document")

            table = DocxTable(child, doc)
            matrix = [
                [cell.text.strip() for cell in row.cells]
                for row in table.rows
            ]
            slug = current["slug"]
            table_counters[slug] = table_counters.get(slug, 0) + 1
            obj = build_table_object(slug, table_counters[slug], matrix)
            if obj:
                current["tables"].append(obj)

    for sec in sections:
        sec["signals"] = compute_density(sec["body_blocks"], len(sec["tables"]))

    return {
        "source_file": path.name,
        "file_type":   "docx",
        "sections":    sections,
    }


# ---------------------------------------------------------------------------
# XLSX extractor
# ---------------------------------------------------------------------------

def extract_xlsx(path: Path) -> Dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)
    sections = []

    for i, sheet in enumerate(wb.worksheets, start=1):
        matrix = []
        for r in sheet.iter_rows():
            matrix.append(
                ["" if c.value is None else str(c.value).strip() for c in r]
            )

        slug  = slugify(sheet.title, f"sheet_{i}")
        table = build_table_object(slug, 1, matrix)

        body_blocks = []
        tables      = [table] if table else []

        sections.append({
            "index":       i,
            "slug":        slug,
            "label":       f"Sheet: {sheet.title}",
            "layout":      "Spreadsheet",
            "title":       sheet.title,
            "body_blocks": body_blocks,
            "tables":      tables,
            "notes":       "",
            "source": {
                "file":     path.name,
                "type":     "xlsx",
                "location": f"Sheet: {sheet.title}",
            },
            "signals": compute_density(body_blocks, len(tables)),
        })

    return {
        "source_file": path.name,
        "file_type":   "xlsx",
        "sections":    sections,
    }


# ---------------------------------------------------------------------------
# PDF extractor
# Primary:  pdfplumber  (text + vector tables)
# Fallback: pytesseract + pdf2image  (OCR for image-heavy pages)
#
# NOTE: OCR only activates on pages where pdfplumber returns no text.
# Requires Tesseract installed separately:
#   Windows: https://github.com/UB-Mannheim/tesseract/wiki
#   Linux:   sudo apt install tesseract-ocr
# ---------------------------------------------------------------------------

def _ocr_available() -> bool:
    try:
        import pytesseract
        import pdf2image  # noqa: F401
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


_OCR_AVAILABLE: Optional[bool] = None   # cached after first check


def _page_ocr(pdf_path: Path, page_number: int) -> str:
    """
    OCR a single PDF page (1-indexed) and return extracted text.
    Only called when pdfplumber finds no text on the page.
    """
    import pytesseract
    from pdf2image import convert_from_path

    images = convert_from_path(
        str(pdf_path),
        first_page=page_number,
        last_page=page_number,
        dpi=300,
    )
    if not images:
        return ""
    return pytesseract.image_to_string(images[0]).strip()


def extract_pdf(path: Path) -> Dict[str, Any]:
    global _OCR_AVAILABLE
    import pdfplumber

    if _OCR_AVAILABLE is None:
        _OCR_AVAILABLE = _ocr_available()
        if not _OCR_AVAILABLE:
            print(
                "[convert] OCR fallback unavailable — install pytesseract + "
                "pdf2image + Tesseract to capture image-embedded text in PDFs."
            )

    sections = []

    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text       = page.extract_text() or ""
            tables_raw = page.extract_tables() or []

            # OCR fallback: only if pdfplumber found no text at all on this page
            ocr_used = False
            if not text.strip() and _OCR_AVAILABLE:
                try:
                    text    = _page_ocr(path, i)
                    ocr_used = True
                except Exception as exc:
                    print(f"[convert] OCR failed on page {i} of {path.name}: {exc}")

            body_blocks = []
            if text.strip():
                cleaned = "\n".join(
                    line.strip()
                    for line in text.splitlines()
                    if line.strip()
                )
                body_blocks.append(cleaned)

            slug   = slugify(f"page_{i}", f"page_{i}")
            tables = []

            for ti, matrix in enumerate(tables_raw, start=1):
                norm_matrix = [
                    [("" if c is None else str(c).strip()) for c in row]
                    for row in matrix
                ]
                obj = build_table_object(slug, ti, norm_matrix)
                if obj:
                    tables.append(obj)

            signals = compute_density(body_blocks, len(tables))

            sections.append({
                "index":       i,
                "slug":        slug,
                "label":       f"Page {i}",
                "layout":      "PDF Page (OCR)" if ocr_used else "PDF Page",
                "title":       f"Page {i}",
                "body_blocks": body_blocks,
                "tables":      tables,
                "notes":       "",
                "source": {
                    "file":     path.name,
                    "type":     "pdf",
                    "location": f"Page {i}",
                },
                "signals": signals,
            })

    return {
        "source_file": path.name,
        "file_type":   "pdf",
        "sections":    sections,
    }


# ---------------------------------------------------------------------------
# TXT / MD extractor
# ---------------------------------------------------------------------------

def extract_text_file(path: Path, ftype: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read().strip()

    body    = [content] if content else []
    signals = compute_density(body, 0)

    return {
        "source_file": path.name,
        "file_type":   ftype,
        "sections": [{
            "index":       1,
            "slug":        slugify(path.stem),
            "label":       "Section 1",
            "layout":      "Text",
            "title":       path.stem,
            "body_blocks": body,
            "tables":      [],
            "notes":       "",
            "source": {
                "file":     path.name,
                "type":     ftype,
                "location": "Entire document",
            },
            "signals": signals,
        }],
    }


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def extract_file(path: Path) -> Dict[str, Any]:
    ext = path.suffix.lower()
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext in (".xlsx", ".xlsm"):
        return extract_xlsx(path)
    if ext == ".txt":
        return extract_text_file(path, "txt")
    if ext == ".md":
        return extract_text_file(path, "md")
    if ext == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported extension: {ext}")


# ---------------------------------------------------------------------------
# File discovery
# ---------------------------------------------------------------------------

def discover_files(input_dir: Path, output_dir: Path) -> List[Path]:
    """
    Recursively find all supported files under input_dir.
    Skips: Office temp files, already-converted outputs, and output_dir itself.
    """
    results = []
    for path in sorted(input_dir.rglob("*"), key=lambda p: str(p).lower()):
        # Don't recurse into the output directory
        try:
            path.relative_to(output_dir)
            continue  # it's inside the output dir, skip
        except ValueError:
            pass

        if not path.is_file():
            continue
        if path.suffix.lower() not in INCLUDE_EXTENSIONS:
            continue
        if any(path.name.startswith(p) for p in SKIP_PREFIXES):
            continue
        if any(path.name.endswith(s) for s in SKIP_SUFFIXES):
            continue

        results.append(path)
    return results


# ---------------------------------------------------------------------------
# Output directory
# ---------------------------------------------------------------------------

def resolve_output_dir(base_output: Path, use_timestamp: bool) -> Path:
    if use_timestamp:
        timestamp = datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")
        out_dir   = base_output / timestamp
    else:
        out_dir = base_output
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Recursively convert supported files under --input to "
            "*_structured.md + *_structured.json in --output."
        )
    )
    parser.add_argument(
        "--input", "-i",
        type=Path,
        default=Path.cwd(),
        help="Root directory to search. Default: current working directory."
    )
    parser.add_argument(
        "--output", "-o",
        type=Path,
        default=None,
        help="Base output directory. Default: <input>/output."
    )
    parser.add_argument(
        "--no-timestamp",
        action="store_true",
        help="Write outputs directly into --output without a timestamped subfolder."
    )
    return parser.parse_args()


def main() -> None:
    args       = parse_args()
    input_dir  = args.input.resolve()
    base_output = (args.output or input_dir / "output").resolve()
    output_dir  = resolve_output_dir(base_output, use_timestamp=not args.no_timestamp)

    print(f"[convert] Input:  {input_dir}")
    print(f"[convert] Output: {output_dir}")

    files = discover_files(input_dir, output_dir)
    if not files:
        print("[convert] No supported files found.")
        return

    print(f"[convert] Found {len(files)} file(s)\n")

    ok  = 0
    err = 0
    for path in files:
        rel = path.relative_to(input_dir)
        print(f"  → {rel}")
        try:
            doc = extract_file(path)

            # Mirror the relative subfolder structure inside the output dir
            out_subdir = output_dir / rel.parent
            out_subdir.mkdir(parents=True, exist_ok=True)

            stem = path.stem
            write_json(doc, out_subdir / f"{stem}_structured.json")
            write_markdown(doc, out_subdir / f"{stem}_structured.md")
            ok += 1
        except Exception as exc:
            print(f"     ERROR: {exc}")
            err += 1

    print(f"\n[convert] Done. {ok} succeeded, {err} failed.")


if __name__ == "__main__":
    main()
