"""
file_to_structured_md.py

File-agnostic context ingestion script.

Outputs (per file):
  - <stem>_structured.json
  - <stem>_structured.md

Supported:
  - .pptx
  - .docx
  - .xlsx / .xlsm
  - .txt / .md

Dependencies:
  pip install python-pptx python-docx openpyxl
"""

import os
import re
import json
from typing import Any, Dict, List, Optional

# -------------------------------------------------------------------
# CONFIG
# -------------------------------------------------------------------
OUTPUT_DIR_NAME = "output"
INPUT_DIR_NAME = ""  # "" = same folder as script

INCLUDE_EXTENSIONS = {
    ".pptx",
    ".docx",
    ".xlsx",
    ".xlsm",
    ".txt",
    ".md",
    ".pdf",
}

# Density thresholds (tweak freely)
LOW_DENSITY_WORDS = 80
HIGH_DENSITY_WORDS = 300

# -------------------------------------------------------------------
# Shared helpers
# -------------------------------------------------------------------
def slugify(text: str, default: str = "section") -> str:
    if not text:
        return default
    slug = text.lower()
    slug = re.sub(r"[^a-z0-9]+", "_", slug)
    slug = re.sub(r"_+", "_", slug)
    return slug.strip("_") or default


def write_json(data: Dict[str, Any], path: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def normalize_cell_for_md(text: Any) -> str:
    if text is None:
        return ""
    parts = [p.strip() for p in str(text).splitlines() if p.strip()]
    return " / ".join(parts)


def render_table_markdown_from_matrix(table_rows: List[List[Any]]) -> str:
    if not table_rows:
        return ""

    max_cols = max(len(r) for r in table_rows)
    rows = [r + [""] * (max_cols - len(r)) for r in table_rows]

    header = rows[0]
    body = rows[1:]

    lines = []
    lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")

    for row in body:
        lines.append("| " + " | ".join(normalize_cell_for_md(c) or " " for c in row) + " |")

    return "\n".join(lines)


def compute_density(body_blocks: List[str], table_count: int) -> Dict[str, Any]:
    words = sum(len(b.split()) for b in body_blocks)
    if words < LOW_DENSITY_WORDS and table_count == 0:
        density = "low"
    elif words > HIGH_DENSITY_WORDS or table_count >= 2:
        density = "high"
    else:
        density = "medium"

    return {
        "word_count": words,
        "table_count": table_count,
        "density": density,
    }


def build_table_object(
    section_slug: str,
    table_index: int,
    matrix: List[List[str]],
) -> Optional[Dict[str, Any]]:
    if not matrix or not any(any(c.strip() for c in r) for r in matrix):
        return None

    max_cols = max(len(r) for r in matrix)
    norm = [r + [""] * (max_cols - len(r)) for r in matrix]

    if any(c.strip() for c in norm[0]):
        headers = norm[0]
        rows = norm[1:]
    else:
        headers = []
        rows = norm

    return {
        "id": f"{section_slug}_table_{table_index}",
        "headers": headers,
        "rows": rows,
        "raw": norm,
    }


# -------------------------------------------------------------------
# Markdown writer
# -------------------------------------------------------------------
def write_markdown(doc: Dict[str, Any], path: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"# Source: {doc['source_file']}\n")
        f.write(f"_File type: {doc['file_type']}_\n\n")
        f.write("---\n\n")

        for sec in doc["sections"]:
            idx = sec["index"]
            title = sec["title"] or f"Section {idx}"

            f.write(f"## SECTION {idx}: {title}\n\n")
            f.write(f"_Slug: {sec['slug']}_\n\n")
            f.write(f"_Layout: {sec['layout']}_\n\n")

            # Trace metadata
            src = sec["source"]
            f.write(
                f"_Source: {src['file']} | {src['type']} | {src['location']}_\n\n"
            )

            # Density
            sig = sec["signals"]
            f.write(
                f"_Signals: {sig['density']} "
                f"(words={sig['word_count']}, tables={sig['table_count']})_\n\n"
            )

            # Body
            for block in sec["body_blocks"]:
                for line in block.splitlines():
                    if line.strip():
                        f.write(f"- {line.strip()}\n")
                f.write("\n")

            # Tables
            for i, table in enumerate(sec["tables"], start=1):
                f.write(f"**Table {i} ({table['id']}):**\n\n")
                matrix = (
                    [table["headers"]] + table["rows"]
                    if table["headers"]
                    else table["raw"]
                )
                f.write(render_table_markdown_from_matrix(matrix) + "\n\n")

            if sec["notes"]:
                f.write("**Notes:**\n\n")
                f.write(sec["notes"] + "\n\n")

            f.write("---\n\n")


# -------------------------------------------------------------------
# PPTX extractor
# -------------------------------------------------------------------
def extract_pptx(path: str) -> Dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(s.shapes)
            else:
                yield s

    def text(shape):
        if not getattr(shape, "has_text_frame", False):
            return ""
        return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()

    def table(shape):
        if not getattr(shape, "has_table", False):
            return None
        return [
            [
                "\n".join(p.text for p in cell.text_frame.paragraphs).strip()
                if cell.text_frame
                else ""
                for cell in row.cells
            ]
            for row in shape.table.rows
        ]

    prs = Presentation(path)
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        title = text(slide.shapes.title) if slide.shapes.title else ""
        slug = slugify(title, f"section_{i}")

        body_blocks = []
        raw_tables = []

        for shape in walk(slide.shapes):
            if slide.shapes.title and shape == slide.shapes.title:
                continue

            t = table(shape)
            if t:
                raw_tables.append(t)
                continue

            txt = text(shape)
            if txt:
                body_blocks.append(txt)

        tables = []
        for ti, matrix in enumerate(raw_tables, start=1):
            obj = build_table_object(slug, ti, matrix)
            if obj:
                tables.append(obj)

        signals = compute_density(body_blocks, len(tables))

        sections.append(
            {
                "index": i,
                "slug": slug,
                "label": f"Slide {i}",
                "layout": slide.slide_layout.name if slide.slide_layout else "",
                "title": title,
                "body_blocks": body_blocks,
                "tables": tables,
                "notes": "",
                "source": {
                    "file": os.path.basename(path),
                    "type": "pptx",
                    "location": f"Slide {i}",
                },
                "signals": signals,
            }
        )

    return {
        "source_file": os.path.basename(path),
        "file_type": "pptx",
        "sections": sections,
    }


# -------------------------------------------------------------------
# DOCX extractor
# -------------------------------------------------------------------
def extract_docx(path: str) -> Dict[str, Any]:
    from docx import Document

    doc = Document(path)
    sections = []
    current = None
    idx = 0

    def start(title, style):
        nonlocal idx
        idx += 1
        sec = {
            "index": idx,
            "slug": slugify(title, f"section_{idx}"),
            "label": f"Section {idx}",
            "layout": style,
            "title": title,
            "body_blocks": [],
            "tables": [],
            "notes": "",
            "source": {
                "file": os.path.basename(path),
                "type": "docx",
                "location": f"Section {idx}",
            },
        }
        sections.append(sec)
        return sec

    for para in doc.paragraphs:
        txt = para.text.strip()
        style = para.style.name if para.style else ""

        if style.startswith("Heading") and txt:
            current = start(txt, style)
        else:
            if not current:
                current = start(
                    os.path.splitext(os.path.basename(path))[0], "Document"
                )
            if txt:
                current["body_blocks"].append(txt)

    # tables → last section
    t_idx = 0
    for table in doc.tables:
        matrix = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        t_idx += 1
        obj = build_table_object(current["slug"], t_idx, matrix)
        if obj:
            current["tables"].append(obj)

    for sec in sections:
        sec["signals"] = compute_density(
            sec["body_blocks"], len(sec["tables"])
        )

    return {
        "source_file": os.path.basename(path),
        "file_type": "docx",
        "sections": sections,
    }


# -------------------------------------------------------------------
# XLSX extractor
# -------------------------------------------------------------------
def extract_xlsx(path: str) -> Dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    sections = []

    for i, sheet in enumerate(wb.worksheets, start=1):
        matrix = []
        for r in sheet.iter_rows():
            matrix.append(
                ["" if c.value is None else str(c.value).strip() for c in r]
            )

        slug = slugify(sheet.title, f"sheet_{i}")
        table = build_table_object(slug, 1, matrix)

        body_blocks = []
        tables = [table] if table else []

        sections.append(
            {
                "index": i,
                "slug": slug,
                "label": f"Sheet: {sheet.title}",
                "layout": "Spreadsheet",
                "title": sheet.title,
                "body_blocks": body_blocks,
                "tables": tables,
                "notes": "",
                "source": {
                    "file": os.path.basename(path),
                    "type": "xlsx",
                    "location": f"Sheet: {sheet.title}",
                },
                "signals": compute_density(body_blocks, len(tables)),
            }
        )

    return {
        "source_file": os.path.basename(path),
        "file_type": "xlsx",
        "sections": sections,
    }
# -------------------------------------------------------------------
# PDF extractor
# -------------------------------------------------------------------
def extract_pdf(path: str) -> Dict[str, Any]:
    import pdfplumber

    sections = []

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            tables_raw = page.extract_tables() or []

            body_blocks = []
            if text.strip():
                # Normalize line breaks similar to PPTX text frames
                cleaned = "\n".join(
                    line.strip()
                    for line in text.splitlines()
                    if line.strip()
                )
                body_blocks.append(cleaned)

            tables = []
            slug = slugify(f"page_{i}", f"page_{i}")

            for ti, matrix in enumerate(tables_raw, start=1):
                # Normalize all cells to strings
                norm_matrix = [
                    [("" if c is None else str(c).strip()) for c in row]
                    for row in matrix
                ]

                obj = build_table_object(slug, ti, norm_matrix)
                if obj:
                    tables.append(obj)

            signals = compute_density(body_blocks, len(tables))

            sections.append(
                {
                    "index": i,
                    "slug": slug,
                    "label": f"Page {i}",
                    "layout": "PDF Page",
                    "title": f"Page {i}",
                    "body_blocks": body_blocks,
                    "tables": tables,
                    "notes": "",
                    "source": {
                        "file": os.path.basename(path),
                        "type": "pdf",
                        "location": f"Page {i}",
                    },
                    "signals": signals,
                }
            )

    return {
        "source_file": os.path.basename(path),
        "file_type": "pdf",
        "sections": sections,
    }


# -------------------------------------------------------------------
# TXT / MD extractor
# -------------------------------------------------------------------
def extract_text_file(path: str, ftype: str) -> Dict[str, Any]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read().strip()

    body = [content] if content else []
    signals = compute_density(body, 0)

    return {
        "source_file": os.path.basename(path),
        "file_type": ftype,
        "sections": [
            {
                "index": 1,
                "slug": slugify(os.path.splitext(os.path.basename(path))[0]),
                "label": "Section 1",
                "layout": "Text",
                "title": os.path.splitext(os.path.basename(path))[0],
                "body_blocks": body,
                "tables": [],
                "notes": "",
                "source": {
                    "file": os.path.basename(path),
                    "type": ftype,
                    "location": "Entire document",
                },
                "signals": signals,
            }
        ],
    }


# -------------------------------------------------------------------
# Dispatcher
# -------------------------------------------------------------------
def extract_file(path: str) -> Dict[str, Any]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext in (".xlsx", ".xlsm"):
        return extract_xlsx(path)
    if ext == ".txt":
        return extract_text_file(path, "txt")
    if ext == ".md":
        return extract_text_file(path, "md")
    if ext == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported file: {ext}")


# -------------------------------------------------------------------
# Main
# -------------------------------------------------------------------
def main():
    base = os.path.dirname(os.path.abspath(__file__))
    input_dir = os.path.join(base, INPUT_DIR_NAME) if INPUT_DIR_NAME else base
    output_dir = os.path.join(base, OUTPUT_DIR_NAME)
    os.makedirs(output_dir, exist_ok=True)

    files = [
        f
        for f in os.listdir(input_dir)
        if os.path.splitext(f)[1].lower() in INCLUDE_EXTENSIONS
        and not f.startswith("~$")
    ]

    if not files:
        print("No supported files found.")
        return

    for fn in files:
        path = os.path.join(input_dir, fn)
        print(f"Processing: {fn}")
        doc = extract_file(path)

        stem = os.path.splitext(fn)[0]
        write_json(doc, os.path.join(output_dir, f"{stem}_structured.json"))
        write_markdown(doc, os.path.join(output_dir, f"{stem}_structured.md"))

    print("\nDone.")


if __name__ == "__main__":
    main()
